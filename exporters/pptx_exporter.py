"""Render an editable, classroom-paced 16:9 teaching deck."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from models.blueprint import CourseBlueprint
from models.slide_deck import Slide, SlideDeck

NAVY = RGBColor(24, 49, 78)
BLUE = RGBColor(42, 105, 164)
SKY = RGBColor(220, 238, 250)
YELLOW = RGBColor(255, 196, 68)
ORANGE = RGBColor(236, 117, 65)
GREEN = RGBColor(43, 143, 103)
INK = RGBColor(35, 41, 49)
MUTED = RGBColor(96, 106, 118)
PAPER = RGBColor(249, 251, 253)
WHITE = RGBColor(255, 255, 255)
CODE_BG = RGBColor(28, 33, 41)
CODE_TEXT = RGBColor(231, 235, 240)
CJK_FONT = "Microsoft YaHei"
CODE_FONT = "Menlo"


def _set_background(slide: PptxSlide, color: RGBColor = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _shape(
    slide: PptxSlide,
    kind: MSO_SHAPE,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    *,
    radius: bool = False,
) -> object:
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind
    shape = slide.shapes.add_shape(
        shape_kind, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _add_text(
    slide: PptxSlide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: RGBColor = INK,
    bold: bool = False,
    font: str = CJK_FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    margin: float = 0.06,
    vertical: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
) -> object:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _add_header(slide: PptxSlide, item: Slide, color: RGBColor = BLUE) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.52, 0.45, 0.12, 0.72, color)
    title_size = 30 if len(item.title) > 15 else 34
    _add_text(
        slide,
        item.title,
        0.78,
        0.28,
        11.7,
        0.98,
        size=title_size,
        color=NAVY,
        bold=True,
    )
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.78, 1.27, 11.7, 0.025, SKY)


def _add_footer(slide: PptxSlide, item: Slide) -> None:
    source = " / ".join(item.source_step_ids) if item.source_step_ids else "COURSE"
    _add_text(slide, f"LessonCraft AI · {source}", 0.78, 7.05, 4.4, 0.22, size=9, color=MUTED)
    _add_text(slide, item.id, 11.45, 7.05, 1.0, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_bullet_list(
    slide: PptxSlide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 22,
    marker_color: RGBColor = BLUE,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    for index, text in enumerate(items[:5]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"●  {text}"
        paragraph.font.name = CJK_FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(14)
        if paragraph.runs:
            paragraph.runs[0].font.color.rgb = marker_color


def _add_interaction(slide: PptxSlide, item: Slide) -> None:
    interaction = item.interaction
    if not interaction or interaction.type == "none" or not interaction.prompt:
        return
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 5.94, 11.55, 0.78, YELLOW, radius=True)
    label = {
        "question": "想一想",
        "prediction": "先预测",
        "choice": "选一选",
        "fill_blank": "填一填",
        "debug": "找一找",
        "hands_on": "动手做",
        "reflection": "说一说",
    }.get(interaction.type, "课堂互动")
    _add_text(slide, label, 1.02, 6.08, 1.25, 0.42, size=16, color=NAVY, bold=True)
    _add_text(slide, interaction.prompt, 2.28, 6.01, 9.75, 0.55, size=20, color=NAVY, bold=True)


def _render_title(slide: PptxSlide, item: Slide) -> None:
    _set_background(slide, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.18, YELLOW)
    _add_text(slide, item.title, 1.05, 1.55, 11.2, 1.45, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "\n".join(item.content), 1.55, 3.2, 10.2, 1.35, size=23, color=SKY, align=PP_ALIGN.CENTER)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.2, 5.4, 6.95, 0.64, BLUE, radius=True)
    _add_text(slide, "从问题出发  ·  用代码验证  ·  完成作品", 3.42, 5.5, 6.5, 0.4, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def _render_section(slide: PptxSlide, item: Slide) -> None:
    _set_background(slide, NAVY)
    _add_text(slide, "COURSE CHALLENGE", 0.95, 0.72, 5.4, 0.42, size=15, color=YELLOW, bold=True)
    _add_text(slide, item.title, 0.95, 1.3, 11.0, 1.25, size=43, color=WHITE, bold=True)
    for index, text in enumerate(item.content[:3]):
        left = 0.95 + index * 4.03
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 3.25, 3.62, 1.62, BLUE if index != 1 else GREEN, radius=True)
        _add_text(slide, f"{index + 1:02d}", left + 0.22, 3.45, 0.62, 0.5, size=20, color=YELLOW, bold=True)
        _add_text(slide, text, left + 0.25, 3.95, 3.12, 0.62, size=21, color=WHITE, bold=True)
    _add_text(slide, item.id, 11.45, 7.02, 1.0, 0.24, size=9, color=SKY, align=PP_ALIGN.RIGHT)


def _render_question(slide: PptxSlide, item: Slide) -> None:
    _set_background(slide, SKY)
    _add_text(slide, "?", 0.82, 0.55, 1.1, 1.1, size=54, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, item.title, 1.82, 0.52, 10.4, 1.05, size=35, color=NAVY, bold=True)
    prompt = item.interaction.prompt if item.interaction and item.interaction.prompt else (item.content[0] if item.content else "")
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 1.88, 11.32, 2.05, WHITE, radius=True)
    _add_text(slide, prompt, 1.45, 2.15, 10.4, 1.45, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    remaining = [text for text in item.content if text != prompt]
    _add_bullet_list(slide, remaining[:3], 1.35, 4.28, 10.6, 1.25, size=20, marker_color=ORANGE)
    _add_footer(slide, item)


def _render_concept(slide: PptxSlide, item: Slide) -> None:
    _add_header(slide, item, GREEN if item.slide_type == "process" else BLUE)
    items = item.content[:5]
    if item.slide_type == "process":
        for index, text in enumerate(items[:3]):
            left = 0.8 + index * 4.12
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 2.05, 3.45, 2.4, WHITE, radius=True)
            _shape(slide, MSO_SHAPE.OVAL, left + 1.3, 1.67, 0.82, 0.82, BLUE)
            _add_text(slide, str(index + 1), left + 1.3, 1.72, 0.82, 0.68, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, text, left + 0.25, 2.45, 2.95, 1.25, size=23, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            if index < 2:
                _add_text(slide, "→", left + 3.5, 2.65, 0.55, 0.72, size=31, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    else:
        primary = items[0] if items else ""
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 1.75, 5.2, 3.72, NAVY, radius=True)
        _add_text(slide, primary, 1.28, 2.25, 4.35, 2.65, size=29, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_bullet_list(slide, items[1:], 6.45, 1.9, 5.65, 3.55, size=21, marker_color=GREEN)
    _add_interaction(slide, item)
    _add_footer(slide, item)


def _render_activity(slide: PptxSlide, item: Slide) -> None:
    _add_header(slide, item, ORANGE)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 1.62, 3.3, 3.95, ORANGE, radius=True)
    _add_text(slide, "DO", 1.15, 1.92, 2.7, 0.78, size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, item.interaction.prompt if item.interaction else "完成课堂任务", 1.17, 2.78, 2.65, 2.25, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for index, text in enumerate(item.content[:4]):
        top = 1.72 + index * 0.92
        _shape(slide, MSO_SHAPE.OVAL, 4.62, top, 0.58, 0.58, NAVY)
        _add_text(slide, str(index + 1), 4.62, top + 0.02, 0.58, 0.48, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, text, 5.42, top - 0.05, 6.35, 0.7, size=21, color=INK, bold=index == 0)
    _add_footer(slide, item)


def _render_comparison(slide: PptxSlide, item: Slide) -> None:
    _add_header(slide, item, BLUE)
    midpoint = max(1, (len(item.content) + 1) // 2)
    groups = [item.content[:midpoint], item.content[midpoint:]]
    headings = ["已经会用", "还需解决"]
    fills = [SKY, RGBColor(255, 237, 221)]
    for index in range(2):
        left = 0.85 + index * 6.05
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 1.75, 5.55, 3.75, fills[index], radius=True)
        _add_text(slide, headings[index], left + 0.35, 1.98, 4.8, 0.55, size=22, color=NAVY, bold=True)
        _add_bullet_list(slide, groups[index], left + 0.35, 2.66, 4.85, 2.25, size=20, marker_color=GREEN if index == 0 else ORANGE)
    _add_interaction(slide, item)
    _add_footer(slide, item)


def _render_code(
    slide: PptxSlide, item: Slide, blueprint: CourseBlueprint
) -> None:
    _add_header(slide, item, GREEN)
    code_id = item.code_display.code_example_id if item.code_display else item.code_example_id
    example = next((value for value in blueprint.code_examples if value.id == code_id), None)
    if example is None:
        _render_concept(slide, item)
        return
    highlights = set(item.code_display.highlight_lines if item.code_display else [])
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.78, 1.5, 8.3, 4.35, CODE_BG, radius=True)
    code_box = slide.shapes.add_textbox(Inches(1.04), Inches(1.75), Inches(7.82), Inches(3.85))
    frame = code_box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = Inches(0.04)
    for index, line in enumerate(example.code.splitlines(), start=1):
        paragraph = frame.paragraphs[0] if index == 1 else frame.add_paragraph()
        paragraph.text = f"{index:>2}  {line}"
        paragraph.font.name = CODE_FONT
        paragraph.font.size = Pt(15 if len(example.code.splitlines()) > 9 else 16)
        paragraph.font.color.rgb = YELLOW if index in highlights else CODE_TEXT
        paragraph.font.bold = index in highlights
        paragraph.space_after = Pt(2)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.4, 1.5, 2.95, 4.35, SKY, radius=True)
    _add_text(slide, f"代码来源\n{example.id}", 9.78, 1.83, 2.2, 0.92, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _add_bullet_list(slide, item.content[:3], 9.68, 2.92, 2.35, 2.45, size=15, marker_color=GREEN)
    _add_interaction(slide, item)
    _add_footer(slide, item)


def _render_summary(slide: PptxSlide, item: Slide) -> None:
    _set_background(slide, NAVY)
    _add_text(slide, item.title, 0.9, 0.62, 11.4, 0.9, size=38, color=WHITE, bold=True)
    for index, text in enumerate(item.content[:3]):
        top = 1.8 + index * 1.25
        _shape(slide, MSO_SHAPE.OVAL, 1.0, top, 0.78, 0.78, YELLOW)
        _add_text(slide, str(index + 1), 1.0, top + 0.03, 0.78, 0.65, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, text, 2.05, top - 0.02, 9.65, 0.82, size=24, color=WHITE, bold=True)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 5.68, 11.3, 0.75, BLUE, radius=True)
    prompt = item.interaction.prompt if item.interaction else "用一句话说出今天的核心结论。"
    _add_text(slide, prompt, 1.35, 5.83, 10.6, 0.42, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, item.id, 11.4, 7.02, 1.0, 0.22, size=9, color=SKY, align=PP_ALIGN.RIGHT)


def _render_assignment(slide: PptxSlide, item: Slide) -> None:
    _add_header(slide, item, ORANGE)
    task = item.content[0] if item.content else "完成作品升级"
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 1.65, 7.0, 3.95, NAVY, radius=True)
    _add_text(slide, "YOUR MISSION", 1.28, 1.98, 4.0, 0.45, size=15, color=YELLOW, bold=True)
    _add_text(slide, task, 1.28, 2.5, 6.18, 2.1, size=27, color=WHITE, bold=True)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.25, 1.65, 4.05, 3.95, SKY, radius=True)
    _add_text(slide, "完成检查", 8.62, 1.98, 3.28, 0.52, size=22, color=NAVY, bold=True)
    _add_bullet_list(slide, item.content[1:], 8.62, 2.75, 3.2, 2.3, size=19, marker_color=GREEN)
    _add_interaction(slide, item)
    _add_footer(slide, item)


def _write_notes(slide: PptxSlide, item: Slide) -> None:
    notes = item.speaker_notes
    text = "\n".join(
        [
            "教师逐页讲解提示",
            f"讲解重点：{notes.explanation or '—'}",
            f"课堂提问：{notes.question or '—'}",
            f"演示动作：{notes.demo or '—'}",
            f"常见错误：{notes.warning or '—'}",
            f"时间建议：{notes.suggested_minutes} 分钟",
            "时间口径：仅表示本页投屏讲解时间，不含学生操作、巡视和过渡。",
            f"下一页过渡：{notes.transition or '—'}",
            "",
            f"来源：步骤 {', '.join(item.source_step_ids) or '无'}；"
            f"目标 {', '.join(item.objective_ids) or '无'}；"
            f"知识点 {', '.join(item.knowledge_ids) or '无'}；"
            f"活动 {', '.join(item.activity_ids) or '无'}；"
            f"练习 {', '.join(item.exercise_ids) or '无'}",
        ]
    )
    try:
        slide.notes_slide.notes_text_frame.text = text
    except (AttributeError, NotImplementedError):
        # Older python-pptx releases may expose notes as read-only.
        return


def _set_core_properties(prs: PresentationType, blueprint: CourseBlueprint) -> None:
    prs.core_properties.title = blueprint.course.title
    prs.core_properties.subject = "少儿编程课堂课件"
    prs.core_properties.author = "LessonCraft AI"
    prs.core_properties.comments = "全部内容由 Course Blueprint 派生"


def export_slide_deck_pptx(
    blueprint: CourseBlueprint, slide_deck: SlideDeck
) -> bytes:
    """Export a validated SlideDeck using layout-specific classroom templates."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _set_core_properties(prs, blueprint)
    blank_layout = prs.slide_layouts[6]
    renderers = {
        "title": _render_title,
        "section": _render_section,
        "question": _render_question,
        "concept": _render_concept,
        "activity": _render_activity,
        "comparison": _render_comparison,
        "summary": _render_summary,
        "assignment": _render_assignment,
    }
    for item in slide_deck.slides:
        slide = prs.slides.add_slide(blank_layout)
        _set_background(slide)
        if item.layout == "code":
            _render_code(slide, item, blueprint)
        else:
            renderer = renderers.get(item.layout, _render_concept)
            renderer(slide, item)
        _write_notes(slide, item)
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
