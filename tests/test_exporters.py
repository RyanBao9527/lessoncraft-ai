"""Offline exporter smoke tests."""

from __future__ import annotations

from io import BytesIO

from docx import Document
from pptx import Presentation

from exporters.docx_exporter import export_lesson_plan_docx
from exporters.pptx_exporter import export_slide_deck_pptx
from models.blueprint import CourseBlueprint
from models.lesson_plan import LessonPlan
from models.slide_deck import SlideDeck


def test_pptx_file_can_be_generated(
    demo_models: tuple[CourseBlueprint, LessonPlan, SlideDeck],
) -> None:
    blueprint, _, deck = demo_models
    data = export_slide_deck_pptx(blueprint, deck)
    presentation = Presentation(BytesIO(data))
    assert len(data) > 10_000
    assert len(presentation.slides) == len(deck.slides)
    assert presentation.slide_width / presentation.slide_height > 1.7
    assert len({slide.layout for slide in deck.slides}) >= 6
    notes_text = presentation.slides[1].notes_slide.notes_text_frame.text
    assert "时间建议" in notes_text
    assert "仅表示本页投屏讲解时间" in notes_text
    for slide_model, slide in zip(deck.slides, presentation.slides, strict=True):
        if not slide_model.code_example_id:
            continue
        example = next(
            item
            for item in blueprint.code_examples
            if item.id == slide_model.code_example_id
        )
        expected = "\n".join(
            f"{index:>2}  {line}"
            for index, line in enumerate(example.code.splitlines(), start=1)
        )
        assert any(
            getattr(shape, "text", "") == expected for shape in slide.shapes
        )


def test_docx_file_can_be_generated(
    demo_models: tuple[CourseBlueprint, LessonPlan, SlideDeck],
) -> None:
    blueprint, lesson_plan, deck = demo_models
    data = export_lesson_plan_docx(blueprint, lesson_plan, deck)
    document = Document(BytesIO(data))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    table_text = "\n".join(
        cell.text
        for table in document.tables
        for row in table.rows
        for cell in row.cells
    )
    assert len(data) > 10_000
    assert blueprint.course.title in text
    assert "课堂流程表" in text
    assert "对应PPT页" in table_text
    assert "教师逐页讲解提示" not in text + table_text
    assert "教师追加" in table_text
    assert "投屏挑战" in table_text
    assert "第 8～9 页" in table_text
    assert "PPT 备注时间仅表示逐页投屏讲解建议时间" in text
    for item in blueprint.code_examples:
        assert (text + table_text).count(item.code) == 1
